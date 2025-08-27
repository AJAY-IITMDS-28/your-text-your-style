from fastapi import FastAPI, File, UploadFile, Form
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
import tempfile
import os
from pptx import Presentation
import requests

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.post("/api/generate")
async def generate_presentation(
    text: str = Form(...),
    guidance: str = Form(''),
    api_key: str = Form(...),
    template: UploadFile = File(...)
):
    # Save uploaded template
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
        tmp.write(await template.read())
        template_path = tmp.name

    # Parse template
    prs = Presentation(template_path)

    # Call LLM to split text into slides
    llm_payload = {
        "model": "gpt-3.5-turbo",
        "messages": [
            {"role": "system", "content": f"You are a helpful assistant. {guidance}"},
            {"role": "user", "content": f"Split the following text into PowerPoint slides. Respond as a JSON array of objects with 'title', 'content'.\n{text}"}
        ]
    }
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    llm_url = "https://api.openai.com/v1/chat/completions"
    try:
        llm_response = requests.post(llm_url, json=llm_payload, headers=headers)
        llm_response.raise_for_status()
        slides = llm_response.json()['choices'][0]['message']['content']
    except Exception as e:
        os.remove(template_path)
        return {"error": "LLM API call failed."}

    import json
    try:
        slides_data = json.loads(slides)
    except Exception:
        os.remove(template_path)
        return {"error": "LLM response parsing failed."}

    # Generate new presentation
    new_prs = Presentation(template_path)
    # Remove all slides
    for i in range(len(new_prs.slides)-1, -1, -1):
        rId = new_prs.slides._sldIdLst[i].rId
        new_prs.slides._sldIdLst.remove(new_prs.slides._sldIdLst[i])
        new_prs.part.drop_rel(rId)
    # Add slides from LLM
    for slide in slides_data:
        layout = new_prs.slide_layouts[0]
        s = new_prs.slides.add_slide(layout)
        s.shapes.title.text = slide.get('title', '')
        for shape in s.shapes:
            if shape.has_text_frame and shape != s.shapes.title:
                shape.text = slide.get('content', '')
                break
    # Save new presentation
    out_path = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx').name
    new_prs.save(out_path)
    os.remove(template_path)
    def iterfile():
        with open(out_path, mode="rb") as file_like:
            yield from file_like
        os.remove(out_path)
    return StreamingResponse(iterfile(), media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
